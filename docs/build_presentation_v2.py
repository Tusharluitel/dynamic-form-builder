"""
Dynamic Form Builder — Technical Presentation v2
Matches the clean white style of the user's original Presentation.pptx.
Run: python3 docs/build_presentation_v2.py
Output: Presentation.pptx  (replaces the original)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import zipfile, copy, io, os

# ── Source images (extracted from original Presentation.pptx) ──────────────
IMG_OLD_SCHEMA   = '/tmp/pptx_image1.jpeg'   # slide 2 — drawSQL old schema
IMG_NEW_SCHEMA   = '/tmp/pptx_image2.png'    # slide 3 — dbdiagram new schema
IMG_THEME_ADMIN  = '/tmp/pptx_image3.png'    # slide 4 — theme admin screenshot
IMG_THEME_FOLDER = '/tmp/pptx_image4.png'    # slide 5 — theme folder tree
IMG_MODULE_FOLDER= '/tmp/pptx_image5.png'    # slide 7 — module folder tree

# ── Dimensions — match original (widescreen 16:9) ──────────────────────────
SLIDE_W = Emu(18288000)
SLIDE_H = Emu(10287000)

# ── Colour palette ──────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x1D, 0x4E, 0xD8)   # indigo-700
ACCENT_LT  = RGBColor(0xDB, 0xEA, 0xFE)   # indigo-100
MUTED      = RGBColor(0x6B, 0x72, 0x80)   # gray-500
GREEN      = RGBColor(0x05, 0x96, 0x69)   # emerald-600
GREEN_LT   = RGBColor(0xD1, 0xFA, 0xE5)   # emerald-100
AMBER      = RGBColor(0xD9, 0x77, 0x06)   # amber-600
AMBER_LT   = RGBColor(0xFE, 0xF3, 0xC7)   # amber-100
RED        = RGBColor(0xDC, 0x26, 0x26)   # red-600
RED_LT     = RGBColor(0xFE, 0xE2, 0xE2)   # red-100
SLATE      = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100 (light bg for rows)
CODE_BG    = RGBColor(0xF8, 0xFA, 0xFC)   # near-white code bg
CODE_FG    = RGBColor(0x1E, 0x40, 0xAF)   # blue-800 code text

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # fully blank


# ════════════════════════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════════════════════════

def tb(slide, text, l, t, w, h,
       size=16, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return box


def tb_lines(slide, lines, l, t, w, h, default_size=14, default_color=BLACK, wrap=True):
    """lines = list of (text, size, bold, color) or plain strings."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, sz, bold, col = item, default_size, False, default_color
        else:
            txt  = item[0]
            sz   = item[1] if len(item) > 1 else default_size
            bold = item[2] if len(item) > 2 else False
            col  = item[3] if len(item) > 3 else default_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(sz)
        run.font.bold  = bold
        run.font.color.rgb = col
    return box


def rect(slide, l, t, w, h, fill=WHITE, line_color=None, line_pt=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def pill(slide, l, t, w, h, fill, text, text_color=WHITE, size=11, bold=True):
    r = rect(slide, l, t, w, h, fill=fill)
    tb(slide, text, l, t + Inches(0.02), w, h, size=size,
       bold=bold, color=text_color, align=PP_ALIGN.CENTER)
    return r


def h_rule(slide, top, color=ACCENT, thick=2):
    rect(slide, Inches(0.5), top, SLIDE_W - Inches(1.0), Pt(thick), fill=color)


def slide_header(slide, title, subtitle=None, pill_text=None, pill_color=ACCENT):
    """Standard slide header with title, optional subtitle and pill badge."""
    h_rule(slide, Inches(0.48))
    tb(slide, title, Inches(0.5), Inches(0.55), Inches(12.5), Inches(0.55),
       size=28, bold=True, color=BLACK)
    if pill_text:
        pill(slide, Inches(13.3), Inches(0.58), Inches(1.7), Inches(0.4),
             pill_color, pill_text, size=11)
    if subtitle:
        tb(slide, subtitle, Inches(0.5), Inches(1.15), Inches(16.0), Inches(0.4),
           size=14, color=MUTED, italic=True)


def page_num(slide, n, total=20):
    tb(slide, f'{n} / {total}',
       SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.42), Inches(1.2), Inches(0.3),
       size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def code_box(slide, code, l, t, w, h, size=10.5):
    rect(slide, l, t, w, h, fill=CODE_BG,
         line_color=RGBColor(0xCB, 0xD5, 0xE1), line_pt=0.75)
    box = slide.shapes.add_textbox(
        l + Inches(0.18), t + Inches(0.12),
        w - Inches(0.36), h - Inches(0.24))
    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.strip('\n').split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = CODE_FG
        run.font.name  = 'Courier New'


def bullet_table(slide, rows, l, t, w,
                 row_h=Inches(0.44), label_w=Inches(3.5),
                 size=13, alt=True):
    """rows = list of (label, value) or plain strings."""
    for i, row in enumerate(rows):
        bg = SLATE if (alt and i % 2 == 0) else WHITE
        rect(slide, l, t + i * row_h, w, row_h, fill=bg,
             line_color=RGBColor(0xE2, 0xE8, 0xF0), line_pt=0.5)
        if isinstance(row, tuple):
            lbl, val = row
            tb(slide, lbl, l + Inches(0.15), t + i * row_h + Pt(4),
               label_w, row_h, size=size, bold=True, color=ACCENT)
            tb(slide, val, l + label_w + Inches(0.1),
               t + i * row_h + Pt(4),
               w - label_w - Inches(0.2), row_h, size=size, color=BLACK)
        else:
            tb(slide, '▸  ' + row,
               l + Inches(0.15), t + i * row_h + Pt(4),
               w - Inches(0.3), row_h, size=size, color=BLACK)


def callout(slide, text, l, t, w, h,
            bg=ACCENT_LT, text_color=ACCENT, size=13, bold=False, icon='ℹ'):
    rect(slide, l, t, Pt(4), h, fill=ACCENT)  # left accent strip
    rect(slide, l + Pt(4), t, w - Pt(4), h, fill=bg)
    tb(slide, icon + '  ' + text,
       l + Inches(0.15), t + Pt(5), w - Inches(0.2), h,
       size=size, color=text_color, bold=bold)


def add_image(slide, img_path, l, t, w, h):
    slide.shapes.add_picture(img_path, l, t, w, h)


TOTAL = 20


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
rect(s, 0, 0, Inches(0.6), SLIDE_H, fill=ACCENT)
rect(s, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill=ACCENT)

tb(s, "Dynamic Form Builder",
   Inches(1.0), Inches(2.6), Inches(15), Inches(1.2),
   size=48, bold=True, color=BLACK)
tb(s, "Technical Progress Report",
   Inches(1.0), Inches(3.9), Inches(15), Inches(0.65),
   size=28, color=ACCENT, bold=True)
tb(s, "Drupal 7  ·  Custom Module  ·  Custom Theme  ·  Response Module",
   Inches(1.0), Inches(4.75), Inches(15), Inches(0.4),
   size=16, color=MUTED)
page_num(s, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Old Schema
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Old Database Schema", "Initial design — 8 tables, no soft delete, no audit trail", "Before", RED)
add_image(s, IMG_OLD_SCHEMA,
          Inches(0.5), Inches(1.6), Inches(9.5), Inches(7.8))

# Annotations on the right
tb(s, "What the old schema had:", Inches(10.3), Inches(1.65), Inches(7.2), Inches(0.4),
   size=15, bold=True, color=BLACK)
bullet_table(s, [
    "forms, sections, questions, question_options",
    "answers, answer_file, responses",
    "form_members (users → forms)",
    "validation_rules stored as JSON blob in questions",
    "No deleted_at column — hard deletes only",
    "No audit log table",
    "No invitations table",
    "No tags tables",
    "No response_sessions table",
], Inches(10.3), Inches(2.1), Inches(7.4),
   row_h=Inches(0.44), size=13, label_w=Inches(0.01))

callout(s, "datetime columns used instead of Unix timestamps — changed in new schema",
        Inches(10.3), Inches(6.2), Inches(7.4), Inches(0.55),
        bg=AMBER_LT, text_color=AMBER, icon='⚠')
page_num(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — New Schema
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "New Database Schema", "14 tables — all with soft delete; added via hook_schema() + update hooks", "After", GREEN)
add_image(s, IMG_NEW_SCHEMA,
          Inches(0.5), Inches(1.6), Inches(9.3), Inches(7.8))

tb(s, "What changed:", Inches(10.1), Inches(1.65), Inches(7.6), Inches(0.4),
   size=15, bold=True, color=BLACK)
bullet_table(s, [
    ("+  deleted_at / deleted_by",  "Added to ALL 5 entity tables — enables soft delete"),
    ("+  dynamic_form_audit_log",   "Append-only log; old_value + new_value as JSON"),
    ("+  dynamic_form_invitations", "Token-based email invites; 7-day expiry"),
    ("+  dynamic_form_tags",        "Canonical tag dictionary — normalised values"),
    ("+  dynamic_form_tag_answers", "Junction: response × question × tag"),
    ("+  response_sessions",        "IP, user_agent — GDPR-isolated table"),
    ("+  question_validations",     "Normalised rules (was JSON blob)"),
    ("~  forms.slug",               "Auto-generated, unique, collision-safe"),
    ("~  forms.visibility",         "p / r / m  (public, restricted, members)"),
    ("~  timestamps",               "Switched from datetime → Unix int"),
], Inches(10.1), Inches(2.1), Inches(7.7),
   row_h=Inches(0.43), size=12, label_w=Inches(3.2))

callout(s, "Schema evolved in 3 update hooks: 7001 (tags), 7002 (width column), 7003 (invitations)",
        Inches(10.1), Inches(6.55), Inches(7.7), Inches(0.55),
        bg=ACCENT_LT, text_color=ACCENT, icon='⚙')
page_num(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Custom Theme overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Custom Theme: dynamic_form",
             "Activated programmatically on module enable — overrides Drupal's default Bartik theme")
add_image(s, IMG_THEME_ADMIN,
          Inches(0.5), Inches(1.6), Inches(8.5), Inches(4.2))

tb(s, "How theme activation works:", Inches(9.3), Inches(1.65), Inches(8.5), Inches(0.4),
   size=15, bold=True, color=BLACK)
code_box(s,
"""// In dynamic_form.install — hook_enable():
theme_enable(array('dynamic_form'));
variable_set('theme_default', 'dynamic_form');
drupal_theme_rebuild();

// hook_disable() — reverts cleanly:
variable_set('theme_default', 'bartik');
theme_disable(array('dynamic_form'));""",
    Inches(9.3), Inches(2.1), Inches(8.5), Inches(2.1))

tb(s, "Key overriding points:", Inches(9.3), Inches(4.35), Inches(8.5), Inches(0.4),
   size=15, bold=True, color=BLACK)
bullet_table(s, [
    ("hook_theme()",    "Registers 'dynamic_form_dashboard' template with sidebar + content vars"),
    ("template.php",    "Preprocesses variables before tpl.php renders (adds user context, nav)"),
    ("page.tpl.php",    "Replaces Drupal's default page HTML — controls nav, header, footer"),
    ("dashboard.tpl.php", "Custom layout: sidebar nav + main content region"),
], Inches(9.3), Inches(4.78), Inches(8.5),
   row_h=Inches(0.43), size=13, label_w=Inches(2.8))
page_num(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Theme folder structure + template files
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Theme File Structure & Template Files",
             "themes/dynamic_form/ — css, js, template, info, template.php")
add_image(s, IMG_THEME_FOLDER,
          Inches(0.5), Inches(1.55), Inches(4.8), Inches(5.5))

tb(s, "Template file roles:", Inches(5.6), Inches(1.6), Inches(12.2), Inches(0.4),
   size=15, bold=True, color=BLACK)
bullet_table(s, [
    ("dynamic_form.info",      "Declares theme name, core=7.x, regions; no dependencies needed"),
    ("template.php",           "Implements template_preprocess_*() hooks — injects PHP vars into tpl"),
    ("template/page.tpl.php",  "Full HTML shell: <head>, navigation bar, <main>, footer rendered here"),
    ("template/dashboard.tpl.php", "Renders $sidebar and $content — used by hook_theme() registration"),
    ("css/style.css",          "All theme CSS — dashboard layout, form cards, builder UI"),
    ("js/main.js",             "Theme-level JS — nav toggle, responsive behaviour"),
], Inches(5.6), Inches(2.05), Inches(12.2),
   row_h=Inches(0.47), size=13, label_w=Inches(3.5))

tb(s, "What 'overriding' means in Drupal 7:", Inches(5.6), Inches(5.38), Inches(12.2), Inches(0.38),
   size=14, bold=True, color=BLACK)

callout(s,
    "Drupal 7 looks for page.tpl.php in the active theme. Because dynamic_form is set as "
    "theme_default, Drupal uses our page.tpl.php instead of Bartik's — giving us full control "
    "over every HTML element on every page. dashboard.tpl.php is only invoked for routes that "
    "call theme('dynamic_form_dashboard', …) registered in hook_theme().",
    Inches(5.6), Inches(5.8), Inches(12.2), Inches(1.3),
    bg=ACCENT_LT, text_color=BLACK, size=13, icon='')
page_num(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dynamic Form Module: folder structure + hooks overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Dynamic Form Module — Structure & Hooks",
             "sites/all/modules/dynamic_form/ — split across .info, .install, .module, includes/, js/, css/")
add_image(s, IMG_MODULE_FOLDER,
          Inches(0.5), Inches(1.55), Inches(5.2), Inches(5.0))

tb(s, "Module file responsibilities:", Inches(6.0), Inches(1.6), Inches(11.8), Inches(0.4),
   size=15, bold=True, color=BLACK)
bullet_table(s, [
    ("dynamic_form.info",    "Module metadata — name, description, core=7.x, dependencies"),
    ("dynamic_form.install", "hook_schema(): all 14 table definitions; hook_enable/disable/uninstall; update hooks 7001-7003"),
    ("dynamic_form.module",  "All hooks: hook_init, hook_permission, hook_menu, hook_theme, hook_cron, hook_mail, hook_form_alter"),
    ("includes/*.inc",       "Domain logic loaded on demand via 'file' key in hook_menu() — not always in memory"),
    ("js/*.js",              "builder.js, members.js, invitations.js, dashboard-search.js, dashboard-delete.js, toast.js"),
    ("css/*.css",            "builder.css, preview.css, toast.css, members.css"),
], Inches(6.0), Inches(2.05), Inches(11.8),
   row_h=Inches(0.47), size=13, label_w=Inches(3.2))

callout(s,
    "The 'file' key in hook_menu() means includes/*.inc files are loaded only when their route is visited — "
    "keeping memory footprint low. The module file itself always loads; it only contains hooks and helpers.",
    Inches(6.0), Inches(5.4), Inches(11.8), Inches(0.8),
    bg=ACCENT_LT, text_color=ACCENT, size=13, icon='ℹ')
page_num(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Custom Hooks (module file) — init, permission, theme, form_alter
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Custom Hooks — Part 1",
             "hook_init · hook_permission · hook_theme · hook_form_alter")

# hook_init
tb(s, "hook_init()", Inches(0.5), Inches(1.55), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
code_box(s,
"""// Loads toast CSS+JS only on relevant paths (dashboard/*, form/*, dynamic-form/*)
// Avoids loading toast assets on every Drupal page sitewide.
if (strpos($path, 'dashboard') === 0 || strpos($path, 'form/') === 0 ...) {
    drupal_add_css($module_path . '/css/toast.css');
    drupal_add_js($module_path . '/js/toast.js');
}""",
    Inches(0.5), Inches(1.98), Inches(8.5), Inches(1.55))

# hook_permission
tb(s, "hook_permission()", Inches(0.5), Inches(3.68), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    "access dashboard",
    "view forms  /  create new form",
    "edit own form  /  edit any form",
    "delete own form  /  delete any form",
    "view soft delete  /  restore deleted",
    "view activity log",
], Inches(0.5), Inches(4.1), Inches(8.5),
   row_h=Inches(0.37), size=12, label_w=Inches(0.01))

# hook_theme
tb(s, "hook_theme()", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
code_box(s,
"""return array(
  'dynamic_form_dashboard' => array(
    'variables' => array('sidebar' => '', 'content' => ''),
    'template'  => 'dashboard',
    // Points to themes/dynamic_form/template/dashboard.tpl.php
    'path' => drupal_get_path('theme', 'dynamic_form') . '/template',
  ),
);""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(1.85))

# hook_form_alter
tb(s, "hook_form_alter()", Inches(9.4), Inches(3.98), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
code_box(s,
"""// Overrides Drupal's generic "Sorry, unrecognized username or password"
// error message on login forms with a more user-friendly message.
if ($form_id == 'user_login' || $form_id == 'user_login_block') {
  $form['#validate'][] = 'dynamic_form_custom_login_error_override';
}""",
    Inches(9.4), Inches(4.4), Inches(8.5), Inches(1.5))
page_num(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Custom Hooks Part 2 — hook_menu, hook_cron, hook_mail
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Custom Hooks — Part 2",
             "hook_menu · hook_cron · hook_mail")

# hook_menu
tb(s, "hook_menu()  — 30+ routes across 5 groups",
   Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    ("/login, /register",              "access callback: user_is_anonymous"),
    ("/form/create, /form/%/edit",     "custom access callback — ownership + member role check"),
    ("/dashboard/**",                  "access callback: user_is_logged_in"),
    ("/dashboard/forms/%/builder",     "_dynamic_form_edit_access() — 3-layer check"),
    ("/dynamic-form/ajax/**",          "AJAX endpoints — return JSON commands, no HTML"),
    ("/dynamic-form/invite/accept/%",  "access callback: TRUE — token is the secret"),
], Inches(0.5), Inches(1.98), Inches(8.6),
   row_h=Inches(0.44), size=12, label_w=Inches(3.5))

callout(s, "'file' key on each route means the .inc is only loaded when that path is visited.",
        Inches(0.5), Inches(4.7), Inches(8.6), Inches(0.5),
        bg=ACCENT_LT, text_color=ACCENT, size=12, icon='ℹ')

# hook_cron
tb(s, "hook_cron()  — email queue processing + invitation expiry",
   Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
code_box(s,
"""function dynamic_form_cron() {
  // Process invitation emails (batch of 20)
  _dynamic_form_process_queue_batch(
    'dynamic_form_invite_notify',
    '_dynamic_form_send_invite_notification', 20
  );
  // Process member-added emails (batch of 20)
  _dynamic_form_process_queue_batch(
    'dynamic_form_member_notify',
    '_dynamic_form_send_member_notification', 20
  );
  // Expire stale pending invitations
  db_update('dynamic_form_invitations')
    ->fields(array('status' => 'expired'))
    ->condition('status', 'pending')
    ->condition('expires_at', REQUEST_TIME, '<')
    ->execute();
}""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(3.45))

# hook_mail
tb(s, "hook_mail()  — two email templates",
   Inches(9.4), Inches(5.58), Inches(8.5), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    ("'member_added'",   "Notifies registered user they were added to a form with their role + URL"),
    ("'invitation_sent'","Sends token accept link to unregistered invitees; includes 7-day expiry notice"),
], Inches(9.4), Inches(6.0), Inches(8.5),
   row_h=Inches(0.47), size=13, label_w=Inches(2.8))
page_num(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DFBToast
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "DFBToast — Custom Notification System",
             "Zero-dependency toast utility; PHP AJAX command + JS renderer; jQuery 1.4.4-compatible")

# Left — architecture
tb(s, "Why custom-built?", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    "No third-party library — zero extra HTTP requests",
    "jQuery .bind() / .animate() / .slideDown() used throughout",
    "(.on() would work with jQuery 1.12.4 — kept .bind() for safety)",
    "Converts Drupal's div.messages into toasts on attach",
    "PHP callbacks fire toasts via custom AJAX command",
], Inches(0.5), Inches(1.98), Inches(8.6),
   row_h=Inches(0.44), size=13, label_w=Inches(0.01))

tb(s, "Public JS API (window.DFBToast):", Inches(0.5), Inches(4.2), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""DFBToast.success(msg [, durationMs])  // 4500 ms default
DFBToast.error(msg   [, durationMs])  // 7000 ms default
DFBToast.warning(msg [, durationMs])  // 5000 ms default
DFBToast.info(msg    [, durationMs])  // 4500 ms default""",
    Inches(0.5), Inches(4.62), Inches(8.6), Inches(1.4))

# Right — PHP side
tb(s, "PHP side — AJAX command:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Defined in dynamic_form.module:
function ajax_command_dfb_toast($type, $message) {
  return array(
    'command'    => 'dfbToast',
    'toast_type' => $type,  // success|error|warning|info
    'message'    => $message,
  );
}

// Usage in any AJAX callback:
$commands[] = ajax_command_dfb_toast('success', t('Saved.'));
$commands[] = ajax_command_replace('#wrapper', $html);
return array('#type' => 'ajax', '#commands' => $commands);""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(2.5))

tb(s, "JS registration (toast.js, runs on DOM ready):", Inches(9.4), Inches(4.62), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Extends Drupal's AJAX command protocol:
Drupal.ajax.prototype.commands.dfbToast =
  function (ajax, response) {
    DFBToast[response.toast_type](response.message);
  };

// Also intercepts Drupal div.messages on attach:
Drupal.behaviors.dfbToast = {
  attach: function (context) {
    $('div.messages', context).once('dfb-toast-convert',
      function () { /* map .status/.error/.warning → toast */ });
  }
};""",
    Inches(9.4), Inches(5.05), Inches(8.5), Inches(2.55))
page_num(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Soft Delete Workflow
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Soft Delete Workflow",
             "All deletes are reversible — deleted_at / deleted_by pattern on 5 entity tables")

# Flow diagram (text-based)
tb(s, "Delete flow:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""User clicks Delete
  → JS opens #dfb-delete-confirm-modal  (dashboard-delete.js)
  → User confirms
  → $.ajax POST to /dynamic-form/ajax/delete/{type}/{id}

PHP callback (dynamic_form_ajax_entity_delete):
  → db_update(table)
       ->fields(array(
           'deleted_at' => REQUEST_TIME,
           'deleted_by' => $user->uid,
         ))
       ->condition('id', $id)
       ->execute()
  → _dynamic_form_audit_log(type, id, 'deleted', $old, NULL)
  → ajax_command_remove('#card-{id}')    // removes card from DOM
  → ajax_command_dfb_toast('success', 'Moved to Trash')""",
    Inches(0.5), Inches(1.98), Inches(8.6), Inches(3.35))

tb(s, "All read queries filter:", Inches(0.5), Inches(5.48), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s, "->condition('deleted_at', NULL, 'IS NULL')",
    Inches(0.5), Inches(5.9), Inches(8.6), Inches(0.6))

# Restore flow
tb(s, "Restore flow:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""User visits /dashboard/trash
  → Sees all rows where deleted_at IS NOT NULL
  → Clicks Restore
  → JS opens restore confirmation modal
  → User confirms

PHP callback (dynamic_form_ajax_entity_restore):
  → db_update(table)
       ->fields(array(
           'deleted_at' => NULL,
           'deleted_by' => NULL,
         ))
       ->condition('id', $id)
       ->execute()
  → _dynamic_form_audit_log(type, id, 'restored', $old, $new)
  → ajax_command_remove('#trash-row-{id}')
  → ajax_command_dfb_toast('success', 'Restored')

Permanent delete:
  → db_delete(table)->condition('id', $id)->execute()
  → Requires 'restore deleted' permission""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(5.05))

callout(s, "Permissions: 'view soft delete' to see Trash page; 'restore deleted' to restore or permanently delete.",
        Inches(0.5), Inches(6.72), Inches(8.6), Inches(0.55),
        bg=AMBER_LT, text_color=AMBER, size=12, icon='🔒')
page_num(s, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Audit Log
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Audit Log",
             "Append-only table — full before/after JSON snapshots of every change")

tb(s, "_dynamic_form_audit_log() — called on every write:", Inches(0.5), Inches(1.55), Inches(17.3), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""_dynamic_form_audit_log($entity_type, $entity_id, $action, $old_value = NULL, $new_value = NULL)

// Examples:
_dynamic_form_audit_log('form',     $id, 'created', NULL,    $fields);
_dynamic_form_audit_log('form',     $id, 'updated', $before, $fields);
_dynamic_form_audit_log('section',  $id, 'deleted', $before, NULL);
_dynamic_form_audit_log('question', $id, 'restored',$before, $after);

// Inside the function:
db_insert('dynamic_form_audit_log')->fields(array(
  'entity_type' => $entity_type,  'entity_id'  => $entity_id,
  'action'      => $action,       'actor_uid'  => $user->uid,
  'old_value'   => json_encode($old_value),    // NULL on 'created'
  'new_value'   => json_encode($new_value),    // NULL on 'deleted'
  'ip_address'  => ip_address(),
  'created_at'  => REQUEST_TIME,
))->execute();""",
    Inches(0.5), Inches(1.98), Inches(17.3), Inches(3.3))

tb(s, "Table design:", Inches(0.5), Inches(5.42), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("Append-only",      "Schema comment: 'Rows must never be updated or deleted'"),
    ("Entities covered", "form, section, question, option, validation, response, answer, member"),
    ("Actions tracked",  "created, updated, deleted, restored, published, unpublished, submitted"),
    ("Indexes",          "entity (type+id), actor_uid, action, created_at, entity_time"),
], Inches(0.5), Inches(5.85), Inches(8.6),
   row_h=Inches(0.44), size=13, label_w=Inches(2.6))

tb(s, "Dashboard page:", Inches(9.4), Inches(5.42), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
callout(s,
    "/dashboard/activity-log — requires 'view activity log' permission. "
    "Displays entity type, entity id, action, actor name, timestamp, and "
    "old/new JSON diff for every recorded change in reverse chronological order.",
    Inches(9.4), Inches(5.85), Inches(8.5), Inches(1.15),
    bg=GREEN_LT, text_color=GREEN, size=13, icon='📋')
page_num(s, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — includes: forms.inc
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "includes/dynamic_form.forms.inc",
             "Form CRUD — create, edit, list, delete + members picker + access callbacks")

tb(s, "Key functions:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("dynamic_form_create_form()",      "Drupal Form API: title, description, visibility, closes_at, members"),
    ("dynamic_form_edit_form()",        "Pre-populates same fields; loads existing members via AJAX picker"),
    ("dynamic_form_create_form_submit()","Generates slug, inserts to DB, saves members, logs audit, sends queue notifications"),
    ("dynamic_form_list_page()",        "Lists forms the current user owns or is a member of"),
    ("dynamic_form_delete_confirm()",   "Drupal confirm form — soft deletes on submit"),
    ("_dynamic_form_members_picker_html()", "Renders the AJAX member search + invite HTML block"),
], Inches(0.5), Inches(1.98), Inches(8.6),
   row_h=Inches(0.44), size=12, label_w=Inches(3.8))

tb(s, "jQuery in forms.inc context (members.js):", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Live user search — debounced 300 ms
$search.bind('keyup', function () {
  clearTimeout(searchTimer);
  searchTimer = setTimeout(function () { _doSearch(q); }, 300);
});

// $.ajax GET to /dynamic-form/ajax/members/search?q=...
// Returns JSON array of {uid, name, mail}
// Results rendered as dropdown; click adds to in-memory members[]

// On form submit — members[] serialized to hidden input:
function _sync() {
  $hidden.val(JSON.stringify(members));
}

// All event binding uses .bind() — compatible with all jQuery versions""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(3.05))

callout(s,
    "members.js runs entirely in JS memory — the member list is only persisted to the DB "
    "when the form is saved. The PHP side reads the hidden JSON input and batch-inserts "
    "into dynamic_form_members, then queues notification emails.",
    Inches(9.4), Inches(5.18), Inches(8.5), Inches(1.05),
    bg=ACCENT_LT, text_color=ACCENT, size=13, icon='ℹ')

callout(s,
    "Slug generation: _dynamic_form_generate_slug() — lowercase + hyphens + collision counter "
    "(my-form → my-form-1 → my-form-2) using a while loop against the DB.",
    Inches(0.5), Inches(6.72), Inches(8.6), Inches(0.65),
    bg=GREEN_LT, text_color=GREEN, size=12, icon='🔗')
page_num(s, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — includes: builder.inc (SortableJS + AJAX)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "includes/dynamic_form.builder.inc",
             "Form Builder page — SortableJS drag-drop, Drupal AJAX Forms, modal question editor")

tb(s, "SortableJS — loaded from CDN:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""drupal_add_js(
  'https://cdn.jsdelivr.net/npm/sortablejs@latest/Sortable.min.js',
  'external'
);

// builder.js — _dfbInitSortable():
function _dfbInitSortable(container, handleSel, itemType, endpoint, groupName) {
  new Sortable(container, {
    handle: handleSel,
    group:  groupName,   // allows dragging between sections
    onEnd:  function (evt) {
      // Collect new order from DOM
      $(evt.to).children('.dfb-question-card')
               .each(function (index) { /* read data-id */ });
      // POST to /dynamic-form/ajax/builder/reorder
      $.ajax({ url: endpoint, type: 'POST', data: payload, ... });
    }
  });
}
// Two Sortable instances: sections + questions-within-section""",
    Inches(0.5), Inches(1.98), Inches(8.6), Inches(3.9))

callout(s, "Drag-drop between sections works via SortableJS group option — questions can be moved across sections.",
        Inches(0.5), Inches(6.02), Inches(8.6), Inches(0.55),
        bg=GREEN_LT, text_color=GREEN, size=12, icon='↔')

# Right — AJAX Form system
tb(s, "Drupal AJAX Form — sections form:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Add Section button has #ajax key:
$form['add_section_btn']['#ajax'] = array(
  'callback' => 'dynamic_form_builder_sections_ajax_callback',
  'wrapper'  => 'dfb-sections-wrapper',
);

// Callback returns commands array:
$commands[] = ajax_command_replace('#dfb-sections-wrapper',
                render($form['sections_wrapper']));
$commands[] = ajax_command_invoke('input[name="section_name"]',
                'val', array(''));  // clear input
return array('#type' => 'ajax', '#commands' => $commands);""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(2.35))

tb(s, "Question edit modal — loaded on demand:", Inches(9.4), Inches(4.47), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Click on question card → $.ajax GET:
//   /dynamic-form/ajax/builder/question/{id}/edit-form
// PHP returns ajax_command_html() injecting form HTML into modal
// Modal contains full Drupal AJAX form (add/remove options, validations)
// Save button calls _dfb_question_form_save_ajax_callback():
$commands[] = ajax_command_replace('#question-{id}', $rendered);
$commands[] = ajax_command_invoke('#dfb-edit-question-modal', 'fadeOut', ...);
$commands[] = ajax_command_dfb_toast('success', t('Saved.'));""",
    Inches(9.4), Inches(4.9), Inches(8.5), Inches(2.15))
page_num(s, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — includes: questions.inc
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "includes/dynamic_form.questions.inc",
             "Question CRUD — 12 question types, #states API, add-more options/validations via AJAX")

tb(s, "12 question types:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("text / textarea",    "Single-line and multi-line text input"),
    ("number",             "Numeric input with min/max validation support"),
    ("email / url",        "Auto-validated via question_validations rule_type"),
    ("select / radio / checkbox", "Options stored in dynamic_form_question_options; select supports allow_multiple"),
    ("date / datetime",    "Date picker; datetime-local input"),
    ("file",               "Custom dropzone UI — no Dropzone.js library"),
    ("rating (stars)",     "1–5 star JS interaction in preview.inc"),
    ("scale",              "Button row (1–10); JS click handler"),
    ("richtext",           "TinyMCE loaded from sites/all/libraries/tinymce/"),
    ("tags",               "Select2 free-form tagging — stores to dynamic_form_tags canonical table"),
], Inches(0.5), Inches(1.98), Inches(8.6),
   row_h=Inches(0.44), size=11.5, label_w=Inches(3.2))

# Right — #states and add-more
tb(s, "Drupal #states API — conditional visibility:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Options wrapper only shows for select/radio/checkbox types:
$form['options_wrapper']['#states'] = array(
  'visible' => array(
    ':input[name="type"]' => array(
      array('value' => 'select'),
      array('value' => 'radio'),
      array('value' => 'checkbox'),
    ),
  ),
);
// Drupal generates the JS automatically — no manual jQuery needed""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(2.2))

tb(s, "Add-more options — AJAX within Drupal Form API:", Inches(9.4), Inches(4.32), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// 'Add Option' button triggers AJAX rebuild:
$form['add_option_btn']['#ajax'] = array(
  'callback' => '_dfb_question_options_ajax_callback',
  'wrapper'  => 'dfb-q-options-wrapper',
);
// submit handler increments count in $form_state
// callback re-renders only the options wrapper (not full page)

// Same pattern for 'Add Validation Rule' button
// Validation rules: min_length, max_length, regex, email, url, numeric""",
    Inches(9.4), Inches(4.75), Inches(8.5), Inches(2.3))
page_num(s, 14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — includes: invitations.inc
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "includes/dynamic_form.invitations.inc",
             "Token-based email invites — two code paths; session-based post-registration accept")

tb(s, "Send flow — two code paths:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""dynamic_form_ajax_invitation_send()  [POST from invitations.js]
  │
  ├── user_load_by_mail($email) found?
  │      YES → _dynamic_form_invite_existing_user()
  │             INSERT dynamic_form_members directly
  │             _dynamic_form_queue_member_notifications()
  │
  └── NO  → _dynamic_form_invite_new_email()
             $token = _dynamic_form_generate_invite_token()
                    = bin2hex(drupal_random_bytes(24))  // 48 hex chars
             INSERT dynamic_form_invitations (status=pending, expires=+7d)
             _dynamic_form_queue_invite_notification()
               → DrupalQueue::get('dynamic_form_invite_notify')
                    ->createItem(array('invite_id' => ..., ...))""",
    Inches(0.5), Inches(1.98), Inches(8.6), Inches(3.55))

tb(s, "Accept flow (public URL):", Inches(0.5), Inches(5.67), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""dynamic_form_invite_accept_page($token):
  → Validate: status=pending AND expires_at > NOW
  → Logged in?  YES → _dynamic_form_accept_invitation()
                       INSERT members, SET status=accepted
  → Anonymous?  → $_SESSION['dfb_invite_token'] = $token
                   drupal_goto('login')   // or 'register'
                   // hook_user_insert reads session after registration""",
    Inches(0.5), Inches(6.1), Inches(8.6), Inches(1.98))

# Right — lifecycle
tb(s, "Invitation lifecycle:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("pending",   "Token issued, awaiting acceptance"),
    ("accepted",  "accepted_at timestamp set; user_id FK populated"),
    ("expired",   "Set by hook_cron() when expires_at < REQUEST_TIME"),
    ("revoked",   "revoked_at + revoked_by set by ajax/invitations/%/revoke"),
], Inches(9.4), Inches(1.98), Inches(8.5),
   row_h=Inches(0.44), size=13, label_w=Inches(2.2))

tb(s, "Resend endpoint:", Inches(9.4), Inches(3.8), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""dynamic_form_ajax_invitation_resend($id):
  → UPDATE invitations SET
      expires_at    = REQUEST_TIME + 7*24*3600,
      resend_count  = resend_count + 1,
      last_resent_at = REQUEST_TIME
  → _dynamic_form_queue_invite_notification() — new email queued""",
    Inches(9.4), Inches(4.22), Inches(8.5), Inches(1.55))

callout(s,
    "invitations.js handles the UI: AJAX send/resend/revoke/list calls, "
    "renders the invitation table with masked emails (_dynamic_form_mask_email()), "
    "status badges, and action buttons.",
    Inches(9.4), Inches(5.92), Inches(8.5), Inches(1.1),
    bg=ACCENT_LT, text_color=ACCENT, size=13, icon='ℹ')
page_num(s, 15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — includes: preview.inc (Select2 + TinyMCE)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "includes/dynamic_form.preview.inc",
             "Form preview renderer — Select2, TinyMCE, custom file dropzone, star/scale interactions")

tb(s, "Select2 — conditional loading:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Only loaded if the form has 'select' or 'tags' questions:
function _dynamic_form_preview_maybe_load_select2($form_id) {
  $has_select2 = db_select('dynamic_form_questions', 'q')
    ->condition('form_id', $form_id)
    ->condition('type', array('select','tags'), 'IN')
    ->condition('deleted_at', NULL, 'IS NULL')
    ->countQuery()->execute()->fetchField();
  if (!$has_select2) { return; }  // skip assets if not needed
  drupal_add_css('sites/all/libraries/select2/select2.css');
  drupal_add_js('sites/all/libraries/select2/select2.min.js', array('weight' => 10));
}

// Select2 init (Drupal behavior):
$('.dfp-select2').select2({ width: '100%', allowClear: TRUE });

// Tags — free-form input using Select2 createSearchChoice:
$('.dfp-tags-select2').select2({
  tags: [],
  ajax: { url: '/dynamic-form/ajax/tags/autocomplete',
          results: function(data) { return { results: data }; } },
  createSearchChoice: function(term, data) {
    return { id: term, text: term };  // allows new tags
  }
});""",
    Inches(0.5), Inches(1.98), Inches(8.6), Inches(4.35))

callout(s,
    "Select2 works because jQuery 1.12.4 is active (jQuery Update module). "
    "Select2 v4 requires .prop() which is only available from jQuery 1.6+.",
    Inches(0.5), Inches(6.47), Inches(8.6), Inches(0.65),
    bg=AMBER_LT, text_color=AMBER, size=12, icon='⚠')

# Right — TinyMCE + file dropzone
tb(s, "TinyMCE — conditional loading:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""function _dynamic_form_preview_maybe_load_tinymce($form_id) {
  // Only loads if form has 'richtext' question
  $tinymce_js = 'sites/all/libraries/tinymce/.../tinymce.min.js';
  drupal_add_js($tinymce_js, array('type' => 'file', 'weight' => 10));
  // Init as inline JS (footer, weight 20):
  tinymce.init({ selector: '.dfp-text-editor-area',
                 toolbar: 'bold italic | link | bullist' });
}""",
    Inches(9.4), Inches(1.98), Inches(8.5), Inches(2.0))

tb(s, "Custom file dropzone (no library):", Inches(9.4), Inches(4.12), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// PHP renders a styled div wrapping a hidden file input:
<div class="dfp-file-dropzone">
  <input type="file" class="dfp-file-hidden" ...>
  <div class="dfp-file-dropzone-inner">Drop file or click</div>
  <div class="dfp-file-name"></div>
</div>

// jQuery shows filename on change:
fileInput.bind('change', function () {
  var name = this.files[0] ? this.files[0].name : '';
  fileInput.closest('.dfp-file-dropzone')
           .find('.dfp-file-name').text(name);
});  // No Dropzone.js library — built from scratch""",
    Inches(9.4), Inches(4.55), Inches(8.5), Inches(2.55))
page_num(s, 16)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — includes: sections.inc, trash.inc, tags.inc, dashboard.inc, auth.inc
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "Remaining Include Files",
             "sections · trash · tags · dashboard · auth")

# sections.inc
tb(s, "dynamic_form.sections.inc", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=ACCENT)
bullet_table(s, [
    "Sections created, renamed, deleted from within the Builder page",
    "Standalone edit page at /dashboard/sections/{id}/edit also exists",
    "Inline rename via builder.js: $.ajax POST to /dynamic-form/ajax/builder/section/{id}/rename",
    "Soft delete + audit log on every write operation",
], Inches(0.5), Inches(1.98), Inches(8.6), row_h=Inches(0.4), size=12, label_w=Inches(0.01))

# trash.inc
tb(s, "dynamic_form.trash.inc", Inches(0.5), Inches(3.72), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=ACCENT)
bullet_table(s, [
    "dynamic_form_dashboard_trash_page() — queries all entities WHERE deleted_at IS NOT NULL",
    "Restore confirmation modal in dashboard-delete.js (same file as delete modal)",
    "AJAX restore: /dynamic-form/ajax/trash/{type}/{id}/restore",
    "Permanent delete: direct db_delete(); requires 'restore deleted' permission",
], Inches(0.5), Inches(4.15), Inches(8.6), row_h=Inches(0.4), size=12, label_w=Inches(0.01))

# tags.inc
tb(s, "dynamic_form.tags.inc", Inches(0.5), Inches(5.88), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=ACCENT)
code_box(s,
"""// Autocomplete endpoint for Select2 ajax on tags question:
dynamic_form_ajax_tags_autocomplete()
  → SELECT standardized_value FROM dynamic_form_tags
      WHERE standardized_value LIKE '%' . $term . '%'
  → Returns JSON array for Select2 results callback

// Tag normalization (module-level helper):
_dynamic_form_normalize_tag($raw)
  → html_entity_decode → mb_strtolower → trim
  → collapse whitespace → strip non-letter/digit/space""",
    Inches(0.5), Inches(6.3), Inches(8.6), Inches(1.87))

# Right — dashboard.inc + auth.inc
tb(s, "dynamic_form.dashboard.inc", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=ACCENT)
bullet_table(s, [
    "dynamic_form_dashboard_page($section) — router for all dashboard views",
    "Renders overview, forms list, sections list, questions list",
    "dynamic_form_activity_log_page() — reads from audit log table",
    "dynamic_form_ajax_entity_delete() — shared soft-delete AJAX endpoint",
    "dashboard-search.js: live client-side filtering of dashboard tables",
], Inches(9.4), Inches(1.98), Inches(8.5), row_h=Inches(0.4), size=12, label_w=Inches(0.01))

tb(s, "dynamic_form.auth.inc", Inches(9.4), Inches(4.1), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=ACCENT)
bullet_table(s, [
    "dynamic_form_login_page() — wraps Drupal's drupal_get_form('user_login') in custom template",
    "dynamic_form_register_page() — wraps user_register_form in custom template",
    "Both routes have access callback: user_is_anonymous — logged-in users are redirected",
    "hook_form_alter() overrides the generic Drupal login error message on both forms",
], Inches(9.4), Inches(4.53), Inches(8.5), row_h=Inches(0.4), size=12, label_w=Inches(0.01))

callout(s,
    "dynamic_form.front.inc handles the public-facing landing page at / — accessible to anonymous users.",
    Inches(9.4), Inches(6.52), Inches(8.5), Inches(0.5),
    bg=ACCENT_LT, text_color=ACCENT, size=12, icon='ℹ')
page_num(s, 17)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — dynamic_form_response module
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "dynamic_form_response Module",
             "Separate sub-module — AJAX save/submit, response dashboard, thank-you page")

tb(s, "Module overview:", Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("dependency",           "Declares dependencies[] = dynamic_form in .info — cannot install standalone"),
    ("own hook_menu()",      "Routes: /dynamic-form-response/ajax/save, /ajax/submit, /dashboard/forms/%/responses, /thankyou/%"),
    ("own hook_permission()", "view form responses, delete form responses, submit form"),
    ("includes/",            "submit.inc — render + save + submit; dashboard.inc — response list + view"),
    ("own JS/CSS",           "Separate assets from main module"),
], Inches(0.5), Inches(1.98), Inches(8.6),
   row_h=Inches(0.47), size=13, label_w=Inches(2.8))

# AJAX save/submit
tb(s, "AJAX save + submit flow:", Inches(0.5), Inches(4.32), Inches(8.6), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""// Autosave (draft) — fires on field change:
POST /dynamic-form-response/ajax/save
  body: { response_id, question_id, value }
  → UPSERT dynamic_form_answers (INSERT or UPDATE)
  → Returns {success: true}  — no full page reload

// Final submit — fires on Submit button:
POST /dynamic-form-response/ajax/submit
  → Validates all required questions server-side
  → UPDATE dynamic_form_responses SET
      status=1, submitted_at=REQUEST_TIME
  → drupal_goto('/thankyou/' . $slug)""",
    Inches(0.5), Inches(4.75), Inches(8.6), Inches(2.88))

# Right
tb(s, "Response dashboard:", Inches(9.4), Inches(1.55), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
bullet_table(s, [
    ("/dashboard/forms/{id}/responses", "Lists all submitted responses for a form"),
    ("/dashboard/forms/{id}/responses/{rid}", "Detail view: renders each question + answer side by side"),
    ("_dynamic_form_response_render_answer_card()", "Handles all question types including file links, tags, richtext"),
    ("Select2 + TinyMCE",               "Loaded conditionally in submit.inc (same as preview.inc pattern)"),
], Inches(9.4), Inches(1.98), Inches(8.5),
   row_h=Inches(0.47), size=13, label_w=Inches(3.5))

tb(s, "response_sessions table:", Inches(9.4), Inches(4.02), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
callout(s,
    "On first load, a row is inserted into dynamic_form_response_sessions "
    "with ip_address, user_agent, and a hashed session_id. This table is "
    "intentionally separate from dynamic_form_responses to isolate PII for GDPR.",
    Inches(9.4), Inches(4.45), Inches(8.5), Inches(1.15),
    bg=RED_LT, text_color=RED, size=13, icon='🔒')

tb(s, "Thank-you page:", Inches(9.4), Inches(5.75), Inches(8.5), Inches(0.38),
   size=14, bold=True, color=BLACK)
code_box(s,
"""dynamic_form_response_thankyou_page($slug)
  → Loads form by slug from dynamic_form_forms
  → Renders confirmation message using form title
  → Accessible anonymously — no permission required""",
    Inches(9.4), Inches(6.18), Inches(8.5), Inches(1.15))
page_num(s, 18)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — jQuery in this project — summary
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
slide_header(s, "jQuery, Select2 & SortableJS — Technical Summary",
             "How and where each library is used across the module")

tb(s, "jQuery 1.12.4  (via jQuery Update module)", Inches(0.5), Inches(1.55), Inches(17.3), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    ("Why upgraded",     "Drupal 7 ships jQuery 1.4.4 — Select2 v4 requires .prop() (jQuery 1.6+)"),
    ("Event binding",    ".bind() used throughout for compatibility safety; .delegate() for dynamic elements"),
    ("AJAX",             "$.ajax() for all custom AJAX calls (member search, invite actions, reorder, section rename)"),
    ("Drupal behaviors", "All JS wrapped in Drupal.behaviors.attach() — safe for AJAX-injected content"),
    (".once()",          "Prevents double-binding when Drupal re-attaches behaviors after AJAX"),
], Inches(0.5), Inches(1.98), Inches(17.3),
   row_h=Inches(0.42), size=13, label_w=Inches(3.0))

tb(s, "Select2  (sites/all/libraries/select2/)", Inches(0.5), Inches(4.18), Inches(17.3), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    ("Used for",         "select questions (multi-select with search) + tags questions (free-form tagging)"),
    ("Loading strategy", "Conditionally loaded: only added to page if form contains select or tags questions"),
    ("Tags mode",        "createSearchChoice callback allows new values; ajax option hits /dynamic-form/ajax/tags/autocomplete"),
    ("Normalization",    "Tag raw input stored in dynamic_form_tag_answers.raw_value; canonical in dynamic_form_tags"),
], Inches(0.5), Inches(4.61), Inches(17.3),
   row_h=Inches(0.42), size=13, label_w=Inches(3.0))

tb(s, "SortableJS  (CDN — jsdelivr)", Inches(0.5), Inches(6.42), Inches(17.3), Inches(0.38),
   size=15, bold=True, color=ACCENT)
bullet_table(s, [
    ("Used for",         "Drag-and-drop reordering of sections and questions in the builder"),
    ("Two instances",    "#dfb-sections-wrapper (sections) + each .dfb-questions-wrapper (questions per section)"),
    ("Cross-section",    "SortableJS group option allows dragging questions between sections"),
    ("Persistence",      "onEnd callback POSTs new order to /dynamic-form/ajax/builder/reorder — PHP updates position column"),
], Inches(0.5), Inches(6.85), Inches(17.3),
   row_h=Inches(0.42), size=13, label_w=Inches(3.0))
page_num(s, 19)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Summary
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
rect(s, 0, 0, Inches(0.6), SLIDE_H, fill=ACCENT)
rect(s, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill=ACCENT)

tb(s, "Summary", Inches(1.0), Inches(0.6), Inches(16), Inches(0.7),
   size=32, bold=True, color=BLACK)
h_rule(s, Inches(1.35), color=ACCENT, thick=2)

cols = [
    ("Schema", [
        "14 custom tables via hook_schema()",
        "Soft delete on all entities",
        "Append-only audit log",
        "3 update hooks (7001–7003)",
        "GDPR-isolated session table",
    ]),
    ("Module Hooks", [
        "hook_init — selective asset loading",
        "hook_menu — 30+ routes",
        "hook_cron — email queue + expiry",
        "hook_mail — 2 email templates",
        "hook_form_alter — login error override",
    ]),
    ("Frontend", [
        "SortableJS — drag-drop builder",
        "Select2 — select + tags inputs",
        "TinyMCE — rich text questions",
        "DFBToast — custom notifications",
        "Custom file dropzone (no library)",
    ]),
    ("Response Module", [
        "Separate installable sub-module",
        "AJAX autosave (draft) + submit",
        "Response dashboard + detail view",
        "Select2 + TinyMCE (conditional)",
        "Thank-you page at /thankyou/{slug}",
    ]),
]

x_positions = [Inches(1.0), Inches(5.3), Inches(9.6), Inches(13.9)]
col_w = Inches(4.0)
for (title, items), x in zip(cols, x_positions):
    rect(s, x, Inches(1.55), col_w, Inches(0.5), fill=ACCENT)
    tb(s, title, x + Inches(0.15), Inches(1.58), col_w - Inches(0.3), Inches(0.45),
       size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    for i, item in enumerate(items):
        bg = SLATE if i % 2 == 0 else WHITE
        rect(s, x, Inches(2.05) + i * Inches(0.54), col_w, Inches(0.54), fill=bg,
             line_color=RGBColor(0xE2, 0xE8, 0xF0), line_pt=0.5)
        tb(s, '▸  ' + item,
           x + Inches(0.12), Inches(2.08) + i * Inches(0.54), col_w - Inches(0.15), Inches(0.5),
           size=12, color=BLACK)

tb(s, "Questions?", Inches(1.0), Inches(5.3), Inches(16), Inches(0.6),
   size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
page_num(s, 20)


# ── Save ────────────────────────────────────────────────────────────────────
OUT = '/Applications/XAMPP/xamppfiles/htdocs/dynamic-form-builder/Presentation.pptx'
prs.save(OUT)
print('Saved:', OUT)
